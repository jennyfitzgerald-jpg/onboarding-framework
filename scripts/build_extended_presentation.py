"""
Deciphex 2026 Kickoff PowerPoint Generator (Extended)
Builds a PPTX from EXTENDED_CONSENSUS_DECK.md with embedded speaker notes.
"""

from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# -----------------------------
# Visual constants (from PDFs)
# -----------------------------

BRAND_ORANGE = RgbColor(0xFF, 0x4D, 0x1C)  # #FF4D1C
DARK_CHARCOAL = RgbColor(0x2D, 0x2D, 0x2D)  # #2D2D2D
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
TEXT_BLACK = RgbColor(0x1A, 0x1A, 0x1A)
GRAY = RgbColor(0xB0, 0xB0, 0xB0)
OFF_WHITE = RgbColor(0xF5, 0xF5, 0xF5)

DECK_TITLE = "Resilience to Results"


# -----------------------------
# Parsing model
# -----------------------------


@dataclass(frozen=True)
class SlideSpec:
    number: int
    title: str
    body_lines: List[str]
    speaker_notes: str


SLIDE_HEADER_RE = re.compile(r"^##\s+SLIDE\s+(\d+):\s*(.+?)\s*$", re.IGNORECASE)


def _strip_md_emphasis(s: str) -> str:
    s = s.strip()
    # Remove leading list bullets
    s = re.sub(r"^\s*[-*]\s+", "", s)
    # Remove basic markdown emphasis around the whole line
    s = re.sub(r"^\*\*(.+)\*\*$", r"\1", s)
    s = re.sub(r"^\*(.+)\*$", r"\1", s)
    s = re.sub(r"^`(.+)`$", r"\1", s)
    return s.strip()


def _clean_line(s: str) -> str:
    s = s.rstrip()
    if s.strip() == "---":
        return ""
    return s


def parse_extended_deck(md_path: Path) -> List[SlideSpec]:
    text = md_path.read_text(encoding="utf-8")
    lines = text.splitlines()

    slides: List[SlideSpec] = []
    cur_num: Optional[int] = None
    cur_title: Optional[str] = None
    cur_lines: List[str] = []

    def flush():
        nonlocal cur_num, cur_title, cur_lines
        if cur_num is None or cur_title is None:
            return
        cleaned = [_clean_line(l) for l in cur_lines]
        cleaned = [l for l in cleaned if l.strip() != ""]

        speaker_notes = ""
        body_lines = cleaned
        try:
            idx = next(i for i, l in enumerate(cleaned) if l.strip().upper() == "**SPEAKER NOTES:**")
            body_lines = cleaned[:idx]
            notes_lines = cleaned[idx + 1 :]
            # Stop notes at first slide separator if it appears in remaining content
            notes_lines = [l for l in notes_lines if l.strip() != "---"]
            speaker_notes = "\n".join([_strip_md_emphasis(l) for l in notes_lines]).strip()
        except StopIteration:
            speaker_notes = ""

        slides.append(
            SlideSpec(
                number=cur_num,
                title=_strip_md_emphasis(cur_title),
                body_lines=[l for l in body_lines if l.strip() != "---"],
                speaker_notes=speaker_notes,
            )
        )

        cur_num, cur_title, cur_lines = None, None, []

    for line in lines:
        m = SLIDE_HEADER_RE.match(line.strip())
        if m:
            flush()
            cur_num = int(m.group(1))
            cur_title = m.group(2).strip()
            cur_lines = []
            continue
        if cur_num is not None:
            cur_lines.append(line)

    flush()
    return slides


# -----------------------------
# Slide type detection
# -----------------------------


def _extract_blockquote(body_lines: List[str]) -> Optional[Tuple[str, str]]:
    # Collect contiguous blockquote lines starting with '>'
    quote_lines = [l.strip()[1:].strip() for l in body_lines if l.strip().startswith(">")]
    if not quote_lines:
        return None
    quote_lines = [l for l in quote_lines if l != ""]

    attribution = ""
    quote_text_lines: List[str] = []
    for l in quote_lines:
        if l.startswith("—") or l.startswith("- —") or l.startswith("--"):
            attribution = l.lstrip("-").strip()
            continue
        quote_text_lines.append(l)

    quote_text = " ".join(quote_text_lines).strip()
    quote_text = quote_text.strip('"').strip()
    quote_text = _strip_md_emphasis(quote_text)

    attribution = attribution.replace("—", "").strip() if attribution else ""
    return (quote_text, attribution)


def _find_first_table(body_lines: List[str]) -> Optional[Tuple[List[str], List[List[str]], List[str]]]:
    # Return headers, rows, remaining lines (without the table)
    table_start = None
    for i, l in enumerate(body_lines):
        if "|" in l and l.strip().startswith("|") and l.strip().endswith("|"):
            # Need a separator row next
            if i + 1 < len(body_lines) and re.match(r"^\|\s*-", body_lines[i + 1].strip()):
                table_start = i
                break
    if table_start is None:
        return None

    table_lines: List[str] = []
    i = table_start
    while i < len(body_lines):
        l = body_lines[i].rstrip()
        if not (l.strip().startswith("|") and "|" in l):
            break
        table_lines.append(l)
        i += 1

    if len(table_lines) < 2:
        return None

    def split_row(row: str) -> List[str]:
        parts = [p.strip() for p in row.strip().strip("|").split("|")]
        return [_strip_md_emphasis(p) for p in parts]

    headers = split_row(table_lines[0])
    rows: List[List[str]] = []
    for row_line in table_lines[2:]:
        if set(row_line.strip()) <= set("|- :"):
            continue
        rows.append(split_row(row_line))

    remaining = body_lines[:table_start] + body_lines[i:]
    remaining = [l for l in remaining if l.strip() != ""]
    return (headers, rows, remaining)


def detect_slide_kind(spec: SlideSpec) -> str:
    title_lower = spec.title.lower()
    if spec.number == 1 and title_lower == "title":
        return "title"
    if "section break" in title_lower:
        return "section_break"
    if _extract_blockquote(spec.body_lines):
        return "quote"
    if _find_first_table(spec.body_lines):
        return "table"
    if any(l.strip().startswith("```") for l in spec.body_lines):
        return "code"
    return "content"


# -----------------------------
# Rendering helpers
# -----------------------------


def _set_bg(slide, color: RgbColor):
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()


def _add_header(slide, title: str):
    # Orange accent bar + gray line like 2025 deck
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(0.5), Inches(1), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_ORANGE
    accent.line.fill.background()

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.75), Inches(0.5), Inches(10.8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.8), Inches(11.8), Inches(0.7))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_BLACK
    p.alignment = PP_ALIGN.LEFT


def _add_footer(slide, slide_number: int, deck_title: str = DECK_TITLE):
    # Left: small orange square + deck title (gray)
    sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.12), Inches(0.18), Inches(0.18))
    sq.fill.solid()
    sq.fill.fore_color.rgb = BRAND_ORANGE
    sq.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(6), Inches(0.3))
    tf = left.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = deck_title
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY

    # Right: "slide" in orange + number in black
    right = slide.shapes.add_textbox(Inches(10.8), Inches(7.02), Inches(2.4), Inches(0.35))
    tf = right.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = "slide"
    p1.font.size = Pt(10)
    p1.font.color.rgb = BRAND_ORANGE
    p1.alignment = PP_ALIGN.RIGHT

    p2 = tf.add_paragraph()
    p2.text = f"{slide_number:02d}"
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_BLACK
    p2.alignment = PP_ALIGN.RIGHT


def _add_orange_triangle(slide):
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(11.8), Inches(6.3), Inches(1.4), Inches(0.95))
    tri.fill.solid()
    tri.fill.fore_color.rgb = BRAND_ORANGE
    tri.line.fill.background()


def add_title_slide(slide_spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_CHARCOAL)

    # Top orange bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_ORANGE
    bar.line.fill.background()

    # Extract title elements from body lines
    lines = [l.strip() for l in slide_spec.body_lines if l.strip() != ""]
    # Expect: **RESILIENCE TO RESULTS**, *The Year...*, presenter line, **mission**
    main = _strip_md_emphasis(lines[0]) if len(lines) > 0 else "RESILIENCE TO RESULTS"
    sub = _strip_md_emphasis(lines[1]) if len(lines) > 1 else "The Year We Turned a Corner"
    presenter = _strip_md_emphasis(lines[2]) if len(lines) > 2 else "Donal O'Shea | CEO | January 2026"
    tagline = _strip_md_emphasis(lines[3]) if len(lines) > 3 else "Precision Pathology Accessible Everywhere"

    t = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(12), Inches(1.2))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = main
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE

    st = slide.shapes.add_textbox(Inches(0.75), Inches(3.55), Inches(12), Inches(0.6))
    tf = st.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = sub
    p.font.size = Pt(30)
    p.font.italic = True
    p.font.color.rgb = BRAND_ORANGE

    pr = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), Inches(8), Inches(0.5))
    tf = pr.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = presenter
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY

    tag = slide.shapes.add_textbox(Inches(6), Inches(6.55), Inches(7.2), Inches(0.5))
    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = BRAND_ORANGE
    p.alignment = PP_ALIGN.RIGHT

    if slide_spec.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_spec.speaker_notes
    return slide


def add_section_break_slide(slide_spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    # Use first H1 in body as headline
    headline = ""
    subtitle = ""
    for l in slide_spec.body_lines:
        if l.strip().startswith("# "):
            headline = _strip_md_emphasis(l.strip()[2:])
            continue
        if l.strip().startswith("*") and l.strip().endswith("*"):
            subtitle = _strip_md_emphasis(l)
            break
    if not headline:
        headline = slide_spec.title

    # Orange accent bar + gray line (like 2025)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.3), Inches(1), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BRAND_ORANGE
    accent.line.fill.background()

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.75), Inches(2.3), Inches(10.8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY
    line.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(12), Inches(1.5))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = headline
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = TEXT_BLACK

    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.75), Inches(4.45), Inches(12), Inches(0.8))
        tf = st.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = GRAY

    _add_orange_triangle(slide)
    _add_footer(slide, slide_spec.number)
    if slide_spec.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_spec.speaker_notes
    return slide


def add_quote_slide(slide_spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)

    quote_data = _extract_blockquote(slide_spec.body_lines) or ("", "")
    quote_text, attribution = quote_data

    # Orange quote marks
    qm1 = slide.shapes.add_textbox(Inches(0.55), Inches(0.55), Inches(1.5), Inches(1.2))
    tf = qm1.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "““"
    p.font.size = Pt(100)
    p.font.bold = True
    p.font.color.rgb = BRAND_ORANGE

    # Separator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.85), Inches(11.8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY
    line.line.fill.background()

    # Quote
    qb = slide.shapes.add_textbox(Inches(0.75), Inches(2.15), Inches(11.8), Inches(2.8))
    tf = qb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = quote_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_BLACK

    # Bottom separator
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(5.25), Inches(11.8), Inches(0.02))
    line2.fill.solid()
    line2.fill.fore_color.rgb = GRAY
    line2.line.fill.background()

    # Attribution (orange)
    ab = slide.shapes.add_textbox(Inches(0.75), Inches(5.5), Inches(8), Inches(0.8))
    tf = ab.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = attribution if attribution else slide_spec.title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRAND_ORANGE

    # Closing quotes
    qm2 = slide.shapes.add_textbox(Inches(11.55), Inches(5.3), Inches(1.5), Inches(1.2))
    tf = qm2.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "””"
    p.font.size = Pt(100)
    p.font.bold = True
    p.font.color.rgb = BRAND_ORANGE

    _add_footer(slide, slide_spec.number)
    if slide_spec.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_spec.speaker_notes
    return slide


def add_content_slide(slide_spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_header(slide, slide_spec.title)

    # Extract bullet points and short paragraphs
    bullets: List[str] = []
    paragraphs: List[str] = []
    for l in slide_spec.body_lines:
        s = l.strip()
        if not s:
            continue
        if s.startswith("[") and s.endswith("]"):
            continue
        if s.startswith("- "):
            bullets.append(_strip_md_emphasis(s))
            continue
        if re.match(r"^\d+\.\s+", s):
            bullets.append(_strip_md_emphasis(s))
            continue
        if s.startswith("# "):
            # Big statement line
            paragraphs.append(_strip_md_emphasis(s[2:]))
            continue
        if s.startswith("**") or s.startswith("*"):
            paragraphs.append(_strip_md_emphasis(s))
            continue
        paragraphs.append(_strip_md_emphasis(s))

    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(11.8), Inches(5.2))
    tf = content_box.text_frame
    tf.clear()
    tf.word_wrap = True

    # Prefer bullets; if none, use paragraphs
    if bullets:
        first = True
        for b in bullets[:12]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = b
            p.font.size = Pt(22)
            p.font.color.rgb = TEXT_BLACK
            p.level = 0
    else:
        first = True
        for para in paragraphs[:6]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = para
            p.font.size = Pt(24)
            p.font.color.rgb = TEXT_BLACK
            p.space_after = Pt(10)

    _add_orange_triangle(slide)
    _add_footer(slide, slide_spec.number)
    if slide_spec.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_spec.speaker_notes
    return slide


def add_table_slide(slide_spec: SlideSpec):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_header(slide, slide_spec.title)

    table_data = _find_first_table(slide_spec.body_lines)
    if not table_data:
        return add_content_slide(slide_spec)
    headers, rows, remaining = table_data

    # Table
    n_cols = max(1, len(headers))
    n_rows = 1 + len(rows)
    table_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(0.75),
        Inches(1.7),
        Inches(11.8),
        Inches(3.8),
    )
    tbl = table_shape.table

    # Headers
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_ORANGE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Rows
    for r, row in enumerate(rows):
        for c in range(n_cols):
            val = row[c] if c < len(row) else ""
            cell = tbl.cell(r + 1, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_BLACK
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER

    # Remaining bullets/notes on slide (below table)
    extras = []
    for l in remaining:
        s = l.strip()
        if not s:
            continue
        if s.startswith("**") or s.startswith("*") or s.startswith("- "):
            extras.append(_strip_md_emphasis(s))
    if extras:
        box = slide.shapes.add_textbox(Inches(0.75), Inches(5.65), Inches(11.8), Inches(1.2))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        first = True
        for e in extras[:6]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = e
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_BLACK

    _add_orange_triangle(slide)
    _add_footer(slide, slide_spec.number)
    if slide_spec.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_spec.speaker_notes
    return slide


def add_code_slide(slide_spec: SlideSpec):
    # Render code block as monospaced text inside a gray box
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_header(slide, slide_spec.title)

    in_code = False
    code_lines: List[str] = []
    bullets: List[str] = []
    for l in slide_spec.body_lines:
        s = l.rstrip("\n")
        if s.strip().startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            code_lines.append(s)
        else:
            if s.strip().startswith("- "):
                bullets.append(_strip_md_emphasis(s))

    if bullets:
        box = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(11.8), Inches(1.2))
        tf = box.text_frame
        tf.clear()
        first = True
        for b in bullets[:4]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = b
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_BLACK

    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.7), Inches(11.8), Inches(4.0))
    rect.fill.solid()
    rect.fill.fore_color.rgb = OFF_WHITE
    rect.line.color.rgb = GRAY

    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.85), Inches(11.3), Inches(3.8))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "\n".join(code_lines[:24])
    p.font.size = Pt(12)
    p.font.name = "Consolas"
    p.font.color.rgb = TEXT_BLACK

    _add_orange_triangle(slide)
    _add_footer(slide, slide_spec.number)
    if slide_spec.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_spec.speaker_notes
    return slide


# -----------------------------
# Main build
# -----------------------------


if __name__ == "__main__":
    root = Path(__file__).resolve().parent
    md_path = root / "EXTENDED_CONSENSUS_DECK.md"
    if not md_path.exists():
        raise FileNotFoundError(f"Missing: {md_path}")

    slides = parse_extended_deck(md_path)
    # Keep only the main 50 slides (ignore appendix)
    slides = [s for s in slides if 1 <= s.number <= 50]
    slides.sort(key=lambda s: s.number)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in slides:
        kind = detect_slide_kind(spec)
        if kind == "title":
            add_title_slide(spec)
        elif kind == "section_break":
            add_section_break_slide(spec)
        elif kind == "quote":
            add_quote_slide(spec)
        elif kind == "table":
            add_table_slide(spec)
        elif kind == "code":
            add_code_slide(spec)
        else:
            add_content_slide(spec)

    out_path = root / "Deciphex_2026_Kickoff_v3.pptx"
    prs.save(str(out_path))
    print(f"[OK] Saved: {out_path}")
    print(f"[OK] Slides: {len(prs.slides)}")

