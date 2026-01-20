"""
Deciphex 2026 Kickoff PowerPoint Generator
Combines consensus deck content with style guide principles
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors (Deciphex Orange theme - from actual PDF slides)
BRAND_PRIMARY = RgbColor(0xFF, 0x4D, 0x1C)  # Deciphex Orange (primary accent)
BRAND_SECONDARY = RgbColor(0xFF, 0x4D, 0x1C)  # Same orange for consistency
BRAND_ACCENT = RgbColor(0xFF, 0x4D, 0x1C)  # Orange accent
DARK_BG = RgbColor(0x2D, 0x2D, 0x2D)  # Dark charcoal (for title slide only)
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RgbColor(0xB0, 0xB0, 0xB0)  # Secondary text gray
OFF_WHITE = RgbColor(0xF5, 0xF5, 0xF5)  # Card backgrounds
DARK_GRAY = RgbColor(0x1A, 0x1A, 0x1A)  # Near black for text
PATHOLYTIX_PINK = RgbColor(0xE9, 0x1E, 0x8C)  # Magenta for sub-brand


def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide"""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def set_shape_fill(shape, color):
    """Set solid fill for a shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle, presenter, tagline):
    """Create the opening title slide with dark theme"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, DARK_BG)
    background.line.fill.background()
    
    # Accent bar at top
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    set_shape_fill(accent_bar, BRAND_SECONDARY)
    accent_bar.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.7), Inches(11.8), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = BRAND_SECONDARY
    p.font.italic = True
    p.alignment = PP_ALIGN.LEFT
    
    # Presenter info
    presenter_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.5), Inches(6), Inches(0.5))
    tf = presenter_box.text_frame
    p = tf.paragraphs[0]
    p.text = presenter
    p.font.size = Pt(20)
    p.font.color.rgb = LIGHT_GRAY
    
    # Tagline bottom right
    tagline_box = slide.shapes.add_textbox(Inches(6), Inches(6.5), Inches(6.5), Inches(0.5))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(18)
    p.font.color.rgb = BRAND_ACCENT
    p.font.italic = True
    p.alignment = PP_ALIGN.RIGHT
    
    return slide


def add_quote_slide(prs, quote, attribution, dark=False):
    """Create an impactful quote slide - 2025 style with white background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background (2025 style - quotes are always on white)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, WHITE)
    background.line.fill.background()
    
    # Opening quote mark - large orange
    quote_mark_open = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1.5), Inches(1.5))
    tf = quote_mark_open.text_frame
    p = tf.paragraphs[0]
    p.text = '\u201C\u201C'  # Double opening quotes
    p.font.size = Pt(100)
    p.font.color.rgb = BRAND_PRIMARY
    p.font.bold = True
    
    # Horizontal line under quote marks
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.8), Inches(11.8), Inches(0.02))
    set_shape_fill(line, LIGHT_GRAY)
    line.line.fill.background()
    
    # Quote text - bold black
    quote_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.8), Inches(2.5))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = quote
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.LEFT
    
    # Horizontal line above attribution
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(5.2), Inches(11.8), Inches(0.02))
    set_shape_fill(line2, LIGHT_GRAY)
    line2.line.fill.background()
    
    # Attribution - orange name
    attr_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.5), Inches(8), Inches(0.8))
    tf = attr_box.text_frame
    p = tf.paragraphs[0]
    p.text = attribution
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRAND_PRIMARY
    p.alignment = PP_ALIGN.LEFT
    
    # Closing quote mark - large orange, bottom right
    quote_mark_close = slide.shapes.add_textbox(Inches(11.5), Inches(5.3), Inches(1.5), Inches(1.5))
    tf = quote_mark_close.text_frame
    p = tf.paragraphs[0]
    p.text = '\u201D\u201D'  # Double closing quotes
    p.font.size = Pt(100)
    p.font.color.rgb = BRAND_PRIMARY
    p.font.bold = True
    
    return slide


def add_section_break_slide(prs, title, subtitle=None):
    """Create a section break slide - 2025 style with white background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background (2025 style)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, WHITE)
    background.line.fill.background()
    
    # Orange accent bar (short) + gray line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.3), Inches(1), Inches(0.06))
    set_shape_fill(accent, BRAND_PRIMARY)
    accent.line.fill.background()
    
    gray_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.75), Inches(2.3), Inches(10.8), Inches(0.02))
    set_shape_fill(gray_line, LIGHT_GRAY)
    gray_line.line.fill.background()
    
    # Section title - black text
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.8), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.italic = True
        p.font.color.rgb = LIGHT_GRAY
    
    # Orange triangle in bottom right (decorative element)
    triangle = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(11.5), Inches(6.2), Inches(1.5), Inches(1))
    set_shape_fill(triangle, BRAND_PRIMARY)
    triangle.line.fill.background()
    
    return slide


def add_content_slide(prs, title, bullet_points, dark=False):
    """Create a content slide with bullets"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg_color = DARK_BG if dark else WHITE
    text_color = WHITE if dark else DARK_GRAY
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, bg_color)
    background.line.fill.background()
    
    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    set_shape_fill(accent, BRAND_SECONDARY)
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    # Bullet points
    bullets_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.8), Inches(5))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = text_color
        p.space_after = Pt(12)
    
    return slide


def add_table_slide(prs, title, headers, rows, dark=False):
    """Create a slide with a data table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg_color = DARK_BG if dark else WHITE
    text_color = WHITE if dark else DARK_GRAY
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, bg_color)
    background.line.fill.background()
    
    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    set_shape_fill(accent, BRAND_SECONDARY)
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table_width = Inches(11.5)
    table_height = Inches(min(0.6 * num_rows, 5))
    col_width = table_width / num_cols
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.9), Inches(1.6), table_width, table_height).table
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if not dark else RgbColor(0x2A, 0x2A, 0x3E)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.color.rgb = text_color
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
    
    return slide


def add_metric_slide(prs, title, metrics, subtitle=None):
    """Create a metrics/KPI slide with big orange numbers - 2025 style"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, WHITE)
    background.line.fill.background()
    
    # Orange accent bar + gray line (2025 style header)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(0.5), Inches(1), Inches(0.06))
    set_shape_fill(accent, BRAND_PRIMARY)
    accent.line.fill.background()
    
    gray_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.75), Inches(0.5), Inches(10.8), Inches(0.02))
    set_shape_fill(gray_line, LIGHT_GRAY)
    gray_line.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.8), Inches(11.8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # Metrics in grid - HUGE orange numbers
    num_metrics = len(metrics)
    cols = min(num_metrics, 4)
    metric_width = Inches(11.8 / cols)
    
    for i, (label, value, growth) in enumerate(metrics):
        col = i % cols
        row = i // cols
        x = Inches(0.75) + (col * metric_width)
        y = Inches(2.0) + (row * Inches(2.5))
        
        # Value (BIG orange number)
        val_box = slide.shapes.add_textbox(x, y, metric_width - Inches(0.2), Inches(1.2))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = BRAND_PRIMARY  # Orange numbers
        p.alignment = PP_ALIGN.CENTER
        
        # Growth indicator (orange)
        if growth:
            growth_box = slide.shapes.add_textbox(x, y + Inches(1.1), metric_width - Inches(0.2), Inches(0.5))
            tf = growth_box.text_frame
            p = tf.paragraphs[0]
            p.text = growth
            p.font.size = Pt(18)
            p.font.color.rgb = BRAND_PRIMARY  # Orange to match
            p.alignment = PP_ALIGN.CENTER
        
        # Label (black)
        label_box = slide.shapes.add_textbox(x, y + Inches(1.5), metric_width - Inches(0.2), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.8), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_two_column_slide(prs, title, left_content, right_content, left_header=None, right_header=None, dark=False):
    """Create a two-column comparison slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    bg_color = DARK_BG if dark else WHITE
    text_color = WHITE if dark else DARK_GRAY
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, bg_color)
    background.line.fill.background()
    
    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    set_shape_fill(accent, BRAND_SECONDARY)
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    
    # Left column header
    if left_header:
        lh_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = lh_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_header
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = BRAND_SECONDARY
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = text_color
        p.space_after = Pt(8)
    
    # Right column header
    if right_header:
        rh_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = rh_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_header
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = BRAND_ACCENT
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = text_color
        p.space_after = Pt(8)
    
    return slide


def add_closing_slide(prs, title, messages, call_to_action):
    """Create a powerful closing slide - 2025 style with white background"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # White background (2025 style)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    set_shape_fill(background, WHITE)
    background.line.fill.background()
    
    # Orange accent bar + gray line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(0.5), Inches(1), Inches(0.06))
    set_shape_fill(accent, BRAND_PRIMARY)
    accent.line.fill.background()
    
    gray_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.75), Inches(0.5), Inches(10.8), Inches(0.02))
    set_shape_fill(gray_line, LIGHT_GRAY)
    gray_line.line.fill.background()
    
    # Title - bold black with key word potentially in orange
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.8), Inches(11.8), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.LEFT
    
    # Messages
    msg_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.8), Inches(3.5))
    tf = msg_box.text_frame
    tf.word_wrap = True
    for i, msg in enumerate(messages):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = msg
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
    
    # Call to action - orange and bold
    cta_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.0), Inches(11.8), Inches(1))
    tf = cta_box.text_frame
    p = tf.paragraphs[0]
    p.text = call_to_action
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRAND_PRIMARY
    p.alignment = PP_ALIGN.LEFT
    
    # Orange triangle decoration bottom right
    triangle = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(11.5), Inches(6.2), Inches(1.5), Inches(1))
    set_shape_fill(triangle, BRAND_PRIMARY)
    triangle.line.fill.background()
    
    return slide


# ============================================================================
# BUILD THE PRESENTATION
# ============================================================================

print("Building Deciphex 2026 Kickoff Presentation...")

# SLIDE 1: Title
slide = add_title_slide(
    prs,
    "RESILIENCE TO RESULTS",
    "The Year We Turned a Corner",
    "Donal O'Shea | CEO | January 2026",
    "Precision Pathology Accessible Everywhere"
)
add_speaker_notes(slide, """Entry energy: Measured, thoughtful—not hype

"Welcome to our 2026 kickoff. This year is different—it's not just a new year, it's a new era for Deciphex."

"We're moving from a year where we proved our resilience to a year where we deliver results."

Wait for room to settle; acknowledge remote participants before beginning.

TRANSITION: "Let me start with a quote that captures where we've been..."
""")

# SLIDE 2: Opening Quote
slide = add_quote_slide(
    prs,
    "Strength does not come from winning. Your struggles develop your strengths. When you go through hardships and decide not to surrender, that is strength.",
    "Arnold Schwarzenegger",
    dark=True
)
add_speaker_notes(slide, """Let the quote sit for 3-4 seconds before speaking.

"I chose this quote deliberately. 2025 tested us in ways we didn't expect. And we didn't surrender."

"The struggles of 2025 didn't weaken us—they revealed who we truly are."

TRANSITION: "Before I talk about the future, I need to acknowledge our present."
""")

# SLIDE 3: Welcome New Joiners
slide = add_content_slide(
    prs,
    "Welcome to the Deciphex Family",
    [
        "~60 new team members since last kickoff",
        "Nearly a third of our company joined during our most challenging year",
        "This presentation is for you as much as anyone",
        "The full context of what we're building together"
    ],
    dark=False
)
add_speaker_notes(slide, """Energy: Inclusive, welcoming

"Look around the room. Look at the faces on your screens."

"Sixty of you have joined us since we last gathered like this—nearly a third of our company."

"You've joined during one of our most challenging years."

"This presentation is for you as much as anyone—the full context of what we're building together."

TRANSITION: "But first, I need to acknowledge what this team went through..."
""")

# SLIDE 4: Acknowledging 2025
slide = add_content_slide(
    prs,
    "A Year of Adversity",
    [
        "We lost Chris in Oxford",
        "Personal tragedies across the team",
        "Health challenges for colleagues and families",
        "Operational headwinds and scaling pains",
        "Summer cashflow pressure—met with €15M Claret Capital venture debt",
        "",
        "We carry these with us. They've shaped us."
    ],
    dark=True
)
add_speaker_notes(slide, """SLOW DOWN significantly here. This is an emotional moment.

"I'm not going to gloss over this. 2025 was hard."

"We lost Chris. For those who knew him—he was one of us. His absence we still feel deeply."

"Beyond this, many of you faced personal challenges—health issues, family struggles."

"I want you to know: I see it. The leadership team sees it. We don't take your resilience for granted."

PAUSE. Let it land.

TRANSITION: "But here's what I've also seen..."
""")

# SLIDE 5: We Turned a Corner
slide = add_table_slide(
    prs,
    "We Turned a Corner",
    ["Achievement", "Impact"],
    [
        ["Northwest London Pathology closed", "Largest NHS contract"],
        ["UK Lab Services launched", "UKAS accredited, operational"],
        ["US digital accessioning live", "Regency, PSL fully digital"],
        ["AI-driven pathologist assignment", "Auto-routing deployed"],
        ["Zero footprint transformation", "Hub elimination advancing"]
    ],
    dark=False
)
add_speaker_notes(slide, """Energy begins to lift here.

"Despite everything, we delivered. Not every target—I'll be honest about that shortly."

"But the foundational pieces—the ones that unlock 2026—they're in place."

"We didn't just survive 2025. We built something."

TRANSITION: "Now let me give you the honest assessment..."
""")

# SLIDE 6: Section Break - The Honest Assessment
slide = add_section_break_slide(
    prs,
    "The Honest Assessment",
    "What worked. What didn't. No spin."
)
add_speaker_notes(slide, """"I believe in radical transparency with this team."

"You can't improve what you won't acknowledge."

TRANSITION: "Let's look at the numbers..."
""")

# SLIDE 7: FY25 By The Numbers
slide = add_metric_slide(
    prs,
    "FY25 By The Numbers",
    [
        ("Group Bookings", "€25.4M", "+52% YoY"),
        ("Group Revenue", "€22.0M", "+45% YoY"),
        ("Clinical Cases", "147,125", "+47% YoY"),
        ("Q4 Bookings", "€11.1M", "44% of annual")
    ],
    subtitle="Strongest Quarter: Q4 — US Revenue: Q1 €38K → Q4 €467K (12x growth)"
)
add_speaker_notes(slide, """"The headline numbers are strong. 52% bookings growth. 45% revenue growth."

"Clinical cases: 147,000 patients whose diagnoses came through our platform."

"But look at the distribution—Q4 was 44% of our year. That's back-loaded. Creates stress and risk."

"2026 needs to be more balanced."

TRANSITION: "Let's look at what specifically worked..."
""")

# SLIDE 8: What Worked
slide = add_table_slide(
    prs,
    "What Worked ✓",
    ["Area", "Achievement"],
    [
        ["Northwest London", "Largest NHS contract closed"],
        ["UK Lab Services", "UKAS accredited, operational"],
        ["Patholytix SaaS", "Strong ARR, high 2026 coverage"],
        ["US Digital Accessioning", "Regency, PSL fully digital"],
        ["Novartis Execution", "3.5M+ slides digitized"],
        ["UK Turnaround Time", "0.8 days vs competitors' 12-28 days"]
    ],
    dark=False
)
add_speaker_notes(slide, """"Let's celebrate what went right."

"Northwest London—years in the making, now executing."

"Lab services—UKAS accreditation unlocks the private market."

"And here's the metric that matters most: 0.8-day turnaround. Competitors? 12-28 days."

"We're not just faster—we're operating in a different dimension of speed."

TRANSITION: "Now the harder conversation..."
""")

# SLIDE 9: What Struggled
slide = add_table_slide(
    prs,
    "What Struggled ✗",
    ["Area", "Reality"],
    [
        ["UK Clinical Growth Velocity", "Below target; young sales org"],
        ["US Scale-Up", "Right pilots, slower progression"],
        ["Canada & ROW", "Lacked bandwidth to execute"],
        ["DAX Analytics Bookings", "Modest; Q1 2026 critical"],
        ["Product-from-Service", "Still struggling with transition"]
    ],
    dark=True
)
add_speaker_notes(slide, """"Now the harder conversation."

"UK clinical—we closed big deals, but velocity of new acquisition was below plan."

"Part of this is a young sales organization learning competitive markets."

"US—we proved the model works, but haven't scaled it yet."

"Canada and ROW—we simply didn't have bandwidth. That's prioritization, not failure."

"DAX Analytics—Georgia is building something real, but bookings were modest."

TRANSITION: "The reality is we're playing in a different league now..."
""")

# SLIDE 10: The Sales Evolution
slide = add_content_slide(
    prs,
    "The Sales Evolution",
    [
        "2023-2024: Won accounts with no backlog provider—lower friction",
        "2025: Entered competitive markets—longer cycles, more sophistication",
        "2026: Target needle-moving accounts: NWLP, Black Country, SynLab, InHealth",
        "",
        "Our turnaround time: 0.8 days vs their 12-28 days",
        "That's 15-35x faster",
        "",
        "We're now playing in a different league."
    ],
    dark=False
)
add_speaker_notes(slide, """"Early on, we won accounts that had no one. The pain was obvious."

"2025, we entered competitive markets. We had to displace incumbents."

"Our turnaround time is 0.8 days vs their 12-28 days. That's 15-35x faster."

"2026, we're targeting accounts that move the needle: 25K+ cases each."

TRANSITION: "Now let me show you what's happening in the market around us..."
""")

# SLIDE 11: Section Break - Market Validation
slide = add_section_break_slide(
    prs,
    "The Market Is Consolidating",
    "And it's validating our thesis."
)
add_speaker_notes(slide, """TRANSITION: "Pay attention to what's happening around us..."
""")

# SLIDE 12: Market Validation
slide = add_two_column_slide(
    prs,
    "Market Validation — M&A Landscape",
    [
        "Tempus acquires Paige AI",
        "$81M (raised $241M)",
        "66% haircut vs funding",
        "Pure AI plays without clinical scale struggle"
    ],
    [
        "AstraZeneca acquires Modella AI",
        "First big pharma acquisition",
        "Strategic assets with clinical integration command premiums"
    ],
    left_header="❌ Weak Position",
    right_header="✓ Strategic Asset",
    dark=True
)
add_speaker_notes(slide, """"Pay attention to what's happening around us."

"Paige was the darling of digital pathology AI. $241M raised. Sold for $81M. 66% loss vs funding."

"Why? AI alone doesn't build a business. You need clinical scale, revenue, operational depth."

"Modella had real clinical integration—AstraZeneca paid a premium."

"Our model—combining clinical services, research platforms, and AI—is what the market now demands."

TRANSITION: "Here's what this means for us..."
""")

# SLIDE 13: What This Means For Us
slide = add_content_slide(
    prs,
    "We Are Building the Strategic Asset",
    [
        "AI-only with no clinical scale → acquired at discounts",
        "AI + clinical operations + operational depth → command premiums",
        "",
        "Our Position:",
        "   • Dual-platform model (clinical + research) validated",
        "   • 9 of Top 10 Pharma already on Patholytix",
        "   • 147K cases processed, data flywheel spinning",
        "",
        "Every case we report, every slide we scan—it compounds."
    ],
    dark=False
)
add_speaker_notes(slide, """"The market is bifurcating."

"If you're AI-only with no revenue, you get acquired at a discount."

"If you're AI plus clinical scale plus operational depth—you're a strategic asset."

"That's what we're building."

TRANSITION: "Now let's talk about how we execute in 2026..."
""")

# SLIDE 14: Section Break - 2026 Strategy
slide = add_section_break_slide(
    prs,
    "2026: The Year of Execution",
    "~50% growth. Clear pathways. No excuses."
)
add_speaker_notes(slide, """Energy: Confident, focused

"Our strategy is simple and focused. No more trying to do everything everywhere."

TRANSITION: "Three pillars..."
""")

# SLIDE 15: Strategic Focus - Three Pillars
slide = add_table_slide(
    prs,
    "Strategic Focus — Three Pillars",
    ["Pillar", "Target", "Key Enabler"],
    [
        ["1. The US Frontier", "3-4x Revenue (100K+ cases)", "MyDiagnexia rollout"],
        ["2. The Fortress", "UK Clinical + SaaS maturity", "NWLP, Black Country, SynLab"],
        ["3. The AI-Native Engine", "10x engineering velocity", "Context Engineering, AI tools"]
    ],
    dark=False
)
add_speaker_notes(slide, """"Our strategy is simple and focused. No more trying to do everything everywhere."

"We turned a corner in 2025. Now we accelerate in the direction we're facing."

"Three pillars: US scale, UK fortress, AI advantage."

NOT DOING: International expansion, new market experiments, scattered innovation.
DOING: Double down on what works.

TRANSITION: "Let me show you the financial targets..."
""")

# SLIDE 16: 2026 Financial Targets
slide = add_table_slide(
    prs,
    "2026 Financial Targets",
    ["Segment", "FY25", "FY26 Target", "Growth"],
    [
        ["Total Revenue", "€22M", "€33M+", "~50%"],
        ["SaaS Research", "€5.3M", "€7.3-7.5M", "+38-42%"],
        ["UK Clinical", "€8.3M", "€18.7M", "+126%"],
        ["US Clinical", "€1.1M", "€3.3M+", "3x"],
        ["Lab Services", "€0.4M", "€3M+", "7.5x"]
    ],
    dark=True
)
add_speaker_notes(slide, """"Here are the numbers. ~50% total growth."

"UK Clinical—this is where NWLP, Black Country, SynLab, InHealth come in."

"US Clinical—3x. This is MyDiagnexia."

"These are ambitious. They're achievable. But they require focus."

Key Dependencies: Northwest London, Black Country, MyDiagnexia US

TRANSITION: "Let me dive into Pillar 1: The US..."
""")

# SLIDE 17: Pillar 1 - The US Frontier
slide = add_content_slide(
    prs,
    "Pillar 1: The US Frontier (MyDiagnexia)",
    [
        "The Insight: American pathologists want empowerment, not replacement",
        "",
        "The Product:",
        "   • 60-second case completion",
        "   • Voice-first 'tap-speak-tap' workflow",
        "   • Self-serve accessioning",
        "   • Revenue enablement for pathologists",
        "",
        "Key Levers: MyDiagnexia Rollout, Web Reporting (Q2), Sonic Expansion, VA Pipeline",
        "",
        "Exit Run Rate Target: 100K+ cases annualized"
    ],
    dark=False
)
add_speaker_notes(slide, """"When we entered the US, we assumed American pathologists wanted what UK pathologists wanted. We were wrong."

"American pathologists are entrepreneurs. They want us to help them do more work, faster, better paid."

"MyDiagnexia is designed for 60-second case completion. Voice-first. Revenue enablement."

"If we prove this in 2026, we've unlocked the largest pathology market in the world."

TRANSITION: "Now the UK Fortress..."
""")

# SLIDE 18: Pillar 2 - The UK Fortress
slide = add_table_slide(
    prs,
    "Pillar 2: The UK Fortress",
    ["Account", "Potential", "Status"],
    [
        ["Northwest London", "50% of H1 new logo slides, £6M+ ARR", "Integration active"],
        ["Black Country", "25K+ cases/year, £1.4M+", "Skin pilot starting"],
        ["SynLab", "25K+ cases/year", "Feb 5-6 exec meeting"],
        ["InHealth", "12K+ cases from August", "Private market entry"]
    ],
    dark=False
)
add_speaker_notes(slide, """"Four accounts define our UK year."

"Northwest London is 50% of H1 new logo slides. Integration can't slip."

"Black Country—£1.4M+ opportunity. Skin pilot is the entry."

"SynLab—Mike Langford meeting with C-suite Feb 5-6. Company-changing partnership potential."

"Our proof point: East Sussex. Before Diagnexia, patients waited 2 months. Now? 2 days."

UK Lab Services: 60K cases target
The Trust Differentiator: 0.8-day average TAT vs competitors' 12-28 days

TRANSITION: "Now the game-changer—our AI transformation..."
""")

# SLIDE 19: Section Break - AI Transformation
slide = add_section_break_slide(
    prs,
    "The AI Transformation",
    "This is the game-changer."
)
add_speaker_notes(slide, """Energy shift here—this should feel different.

"Now I want to talk about something that's going to fundamentally change how we operate."

"This isn't incremental. This is a step change."
""")

# SLIDE 20: The Linux Converter Story
slide = add_two_column_slide(
    prs,
    "The Story That Changed Everything",
    [
        "Traditional Approach (Norbert)",
        "Timeline: 12 months (ongoing)",
        "Cost: ~€70-80K engineering time"
    ],
    [
        "AI-Assisted (Andrea)",
        "Timeline: 24 hours",
        "Cost: ~€100 in tokens"
    ],
    left_header="Old Way",
    right_header="New Way",
    dark=True
)
add_speaker_notes(slide, """"I want to tell you a story."

"Norbert has been working on our Linux converter since January 2025. A year."

"Andrea took the same problem into Claude Code. Done in 24 hours."

"Same outcome. 500x faster. 700x cheaper."

PAUSE for effect.

"This isn't a one-off. This is the future of how we build."

TRANSITION: "This is what we call Context Engineering..."
""")

# SLIDE 21: Context Engineering
slide = add_two_column_slide(
    prs,
    "The New Paradigm — Context Engineering",
    [
        "Weeks/months to ship features",
        "Sprint cycles, story points",
        "Manual implementation",
        "Code-first thinking"
    ],
    [
        "Hours/days to ship features",
        "Continuous delivery",
        "AI orchestration",
        "Problem-first thinking"
    ],
    left_header="Old Model",
    right_header="New Model",
    dark=False
)
add_speaker_notes(slide, """"For decades, we've trained engineers to write code. Line by line."

"That's not the differentiator anymore."

"The differentiator is understanding the problem, decomposing it, orchestrating AI to solve it."

"Our best engineers will be context engineers—not code writers."

The Evolution: Assembly → C → Python → English

"To our competitors still thinking in sprint cycles while we ship features in hours: good luck."

TRANSITION: "Here are the tools we're approving for production..."
""")

# SLIDE 22: AI Tools Approved
slide = add_table_slide(
    prs,
    "AI Tools — Approved for Production",
    ["Tool", "Use Case"],
    [
        ["Cursor IDE", "Comfortable for non-engineering users"],
        ["Claude Code", "Terminal-based, deep reasoning"],
        ["Windsurf", "Alternative IDE experience"]
    ],
    dark=True
)
add_speaker_notes(slide, """"Here's the announcement. Three AI development tools approved for production."

"Marco and InfoSec completed the security assessment."

"The tools are here. The permission is granted. There are no more excuses."

What This Enables:
- BI Dashboards: Days/weeks → Minutes
- API Integrations: Sprint planning → Same day
- Feature Development: Quarterly → Weekly

"This is not a nice-to-have. Our competitors are doing this. If we don't, we fall behind."

TRANSITION: "Our AI team is also transforming how we build AI products..."
""")

# SLIDE 23: The CAV Factory
slide = add_content_slide(
    prs,
    "The CAV Factory — AI Product Vision",
    [
        "Concept Activation Vectors — Rapid AI Deployment",
        "",
        "Old Model: Months to train new endpoint, heavy annotation",
        "CAV Factory: Days to deploy, lightweight concept activation",
        "",
        "Target: New AI endpoints shipping WEEKLY by Q2",
        "",
        "2026 AI Priorities:",
        "   • Company-wide LLM enablement",
        "   • AI-IQA coverage expansion (Derm + GI)",
        "   • AI Risk-Based Assignment: Economy for low-risk, Experience for high-risk",
        "   • Foundation model development"
    ],
    dark=False
)
add_speaker_notes(slide, """"Our AI team has moved from building models to building a CAV Factory."

"Traditional AI: months of annotation, training, validation."

"CAV approach: activate concepts in existing foundation models. Days, not months."

"This allows us to respond to clinical needs in real-time."

"Weekly releases of diagnostic support tools."

TRANSITION: "Now let's talk about our culture and execution model..."
""")

# SLIDE 24: Section Break - Culture & Execution
slide = add_section_break_slide(
    prs,
    "Culture & Execution",
    "Pirates + Navy = Scale"
)
add_speaker_notes(slide, """TRANSITION: "I want to introduce a framework..."
""")

# SLIDE 25: Pirates vs Navy
slide = add_two_column_slide(
    prs,
    "Pirates + Navy — Scaling Without Chaos",
    [
        "Innovation, speed, agility",
        "High-performing individuals",
        "Got us to €22M",
        "Essential for breakthrough thinking"
    ],
    [
        "Process, cadence, predictability",
        "Systematic execution",
        "Gets us to €100M",
        "Enables scale without chaos"
    ],
    left_header="Pirates",
    right_header="Navy",
    dark=True
)
add_speaker_notes(slide, """"I want to introduce a framework: pirates versus navy."

"We've been incredible pirates—innovative, fast, high-performing."

"To handle 300K cases, we need navy discipline: systems, processes, clean handoffs."

"Look at Mike Langford in UK operations. He's a navy operator. And he enabled pirate innovation."

"In 2026, we need more navy. Not to replace pirates—to amplify them."

"What got us to €15M did not get us to €22M. What gets us to €22M will not get us to €100M."

TRANSITION: "Here's what I expect from everyone..."
""")

# SLIDE 26: What We Expect
slide = add_table_slide(
    prs,
    "What We Expect From Everyone",
    ["Expectation", "Meaning"],
    [
        ["Urgency with quality", "Today's problems compound"],
        ["Ownership end-to-end", "Act like owners—because you are"],
        ["Customer obsession", "Their success is our success"],
        ["AI-native ways of working", "No two-speed organization"],
        ['"Good is the enemy of great"', "Don't settle"]
    ],
    dark=False
)
add_speaker_notes(slide, """"Strategy without behaviors is just words."

"I'll be direct: we already have a two-speed problem."

"Some of you have embraced the tools—shipping faster, doing things that weren't possible."

"Others are skeptical. 'AI can't understand our domain.'"

"AI isn't perfect. But it's a 10x multiplier for those who learn to use it."

The AI Adoption Warning:
- Early adopters: 3-5x productivity improvements already
- Skeptics: Still working the old way

"The organization moves together or it fragments."

TRANSITION: "Here's how we double capacity..."
""")

# SLIDE 27: Think 2X
slide = add_content_slide(
    prs,
    "Think 2X — How We Double Capacity",
    [
        "1. Current Organization — Minimize attrition, maximize performance",
        "",
        "2. Lean Processes — Fix the process BEFORE adding people",
        "",
        "3. Smart New People — A-players only; curiosity over current competence",
        "",
        "4. AI Supercharging — Every role, every function, every day",
        "",
        '"Curiosity beats current competence. Always."',
        '"A-players hire A-players. B-players hire C-players."'
    ],
    dark=True
)
add_speaker_notes(slide, """"We don't double by hiring twice as many people."

"We double by getting more from current org, fixing processes before adding headcount, hiring A-players only, and supercharging everything with AI."

"My hiring philosophy: curiosity beats current competence. Always."

TRANSITION: "Before I close, let me remind you why we do this..."
""")

# SLIDE 28: Section Break - The Charge
slide = add_section_break_slide(
    prs,
    "Why We Fight",
    "It's not the numbers. It's the impact."
)
add_speaker_notes(slide, """Energy: Shift to purpose and meaning.

TRANSITION: "Let me remind you why we do this..."
""")

# SLIDE 29: Why We Fight - Impact
slide = add_metric_slide(
    prs,
    "Why We Fight",
    [
        ("Cases Reported FY25", "147,125", "patients served"),
        ("Thousands of Cancers", "Diagnosed", "faster than ever"),
        ("Families", "Informed", "who could plan treatment"),
        ("Drug Development", "Accelerated", "saving future lives")
    ],
    subtitle="This is why we fight through the hard years."
)
add_speaker_notes(slide, """"Before I close—let me remind you why we do this."

"147,000 cases. That's not a metric. That's 147,000 people who needed an answer."

"Thousands of cancers diagnosed. Families who could plan. Treatments that could start."

"That's what your resilience in 2025 produced."

TRANSITION: "Here's where we sit in the landscape..."
""")

# SLIDE 30: The Benchmark
slide = add_table_slide(
    prs,
    "The Benchmark",
    ["Benchmark", "Reality"],
    [
        ["Startups reaching $100M revenue", "0.001-0.002%"],
        ["Venture-backed reaching $100M", "0.2-0.5%"],
        ["Our trajectory", "2027 at current velocity"]
    ],
    dark=True
)
add_speaker_notes(slide, """"Less than half a percent of venture-backed companies reach $100M revenue."

"At our trajectory, we hit that in 2027. Top 0.5%."

"But only if we execute. Only if we don't get comfortable. Only if we stay hungry."

"Good enough is the trap. That's how you end up in the 99.5% that don't make it."

TRANSITION: "Let me close with this..."
""")

# SLIDE 31: Closing Quote
slide = add_quote_slide(
    prs,
    "Good is the enemy of great. And that is one of the key reasons why we have so little that becomes great.",
    "Jim Collins, Good to Great",
    dark=True
)
add_speaker_notes(slide, """Let this sit for 3-4 seconds.

"Don't settle for a good company. Let's build a great one."

TRANSITION: "Here's my call to action..."
""")

# SLIDE 32: The Call to Action
slide = add_closing_slide(
    prs,
    "Resilience to Results",
    [
        "To Our Veterans: Thank you for building what we have. Now help us build what comes next.",
        "",
        "To Our Newcomers: Welcome to the arena. We don't do easy. We do important.",
        "",
        "Three Things I'm Asking:",
        "   1. Embrace the transformation. The AI tools are here. Use them.",
        "   2. Think 2X. Every process. Every system. Every conversation.",
        "   3. Own your outcomes. Act like an owner. Because you are."
    ],
    "Let's make 2026 the year we prove what we're truly capable of."
)
add_speaker_notes(slide, """"2026 is our year. We have everything we need to succeed."

"The capital, the technology, the team, the market position."

"Let's turn ambition into achievement."

PAUSE. Let applause happen naturally.

Post-presentation: Transition to Professor Loughrey keynote.
""")

# SLIDE 33: Final Closing Quote
slide = add_quote_slide(
    prs,
    "The biggest danger is not that our aim is too high and we miss it, but that it is too low and we reach it.",
    "Michelangelo",
    dark=True
)
add_speaker_notes(slide, """Let this resonate.

"Thank you."

End with confidence. Hold the stage. Exit cleanly.
""")

# SLIDE 34: Thank You
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
set_shape_fill(background, DARK_BG)
background.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(11.8), Inches(1.5))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(72)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

tagline_box = slide.shapes.add_textbox(Inches(0.75), Inches(5), Inches(11.8), Inches(0.5))
tf = tagline_box.text_frame
p = tf.paragraphs[0]
p.text = "Precision Pathology Accessible Everywhere"
p.font.size = Pt(24)
p.font.color.rgb = BRAND_ACCENT
p.font.italic = True
p.alignment = PP_ALIGN.CENTER

add_speaker_notes(slide, """Hold for applause.

Transition to next speaker: Professor Loughrey keynote
""")

# Save the presentation
output_path = r"C:\Users\donal.oshea_deciphex\Kickoff workspace\Deciphex_2026_Kickoff_v2.pptx"
prs.save(output_path)
print(f"\n[OK] Presentation saved to: {output_path}")
print(f"[OK] Total slides: {len(prs.slides)}")
print("\nPresentation includes:")
print("  - Dark/light themed slides following style guide")
print("  - Big impact metrics with YoY growth indicators")
print("  - Full speaker notes for every slide")
print("  - Section breaks for emotional arc")
print("  - Quote slides with key inspirational messages")
print("  - Tables for data-heavy content")
print("  - Two-column comparison slides")
print("  - Closing rally sequence")
